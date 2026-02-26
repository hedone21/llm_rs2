#!/usr/bin/env python3
"""Generate D-Bus IPC Architecture PPT — 2 slides, grayscale, formal report style."""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Grayscale Palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_90 = RGBColor(0x2A, 0x2A, 0x2A)   # near black
GRAY_70 = RGBColor(0x4D, 0x4D, 0x4D)   # dark gray
GRAY_50 = RGBColor(0x80, 0x80, 0x80)   # mid gray
GRAY_30 = RGBColor(0xB0, 0xB0, 0xB0)   # light gray
GRAY_15 = RGBColor(0xD9, 0xD9, 0xD9)   # very light gray
GRAY_08 = RGBColor(0xED, 0xED, 0xED)   # near white
GRAY_05 = RGBColor(0xF5, 0xF5, 0xF5)   # barely gray
BORDER = RGBColor(0x99, 0x99, 0x99)     # border gray

FONT_TITLE = "Arial"
FONT_BODY = "Arial"
FONT_MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg_white(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def box(slide, left, top, w, h, fill=GRAY_05, border=BORDER, bw=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = bw
    return s


def rbox(slide, left, top, w, h, fill=GRAY_05, border=BORDER, bw=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = bw
    s.adjustments[0] = 0.06
    return s


def txt(slide, left, top, w, h, text, size=11, color=BLACK, bold=False,
        align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def multiline(slide, left, top, w, h, lines, font=FONT_BODY):
    """lines: list of (text, size, color, bold, align)"""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text, sz, clr, bld = item[0], item[1], item[2], item[3]
        al = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        fn = item[5] if len(item) > 5 else font
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = fn
        p.alignment = al
        p.space_after = Pt(1)
    return tb


def arrow_right(slide, x1, y, x2, color=GRAY_70, width=Pt(1.5)):
    """Horizontal arrow from (x1,y) to (x2,y) with arrowhead."""
    conn = slide.shapes.add_connector(1, x1, y, x2, y)
    conn.line.color.rgb = color
    conn.line.width = width
    # Add arrowhead via XML
    ln = conn.line._ln
    tail_end = ln.find(qn('a:tailEnd'))
    if tail_end is None:
        from lxml import etree
        tail_end = etree.SubElement(ln, qn('a:tailEnd'))
    tail_end.set('type', 'triangle')
    tail_end.set('w', 'med')
    tail_end.set('len', 'med')
    return conn


def add_table(slide, left, top, w, h, rows, cols):
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, h)
    return tbl_shape.table


def style_cell(cell, text, size=9, bold=False, color=BLACK, fill=None,
               align=PP_ALIGN.LEFT, font=FONT_BODY):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # margins
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)
    if fill:
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.find(qn('a:solidFill'))
        if solidFill is not None:
            tcPr.remove(solidFill)
        from lxml import etree
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        srgb = etree.SubElement(sf, qn('a:srgbClr'))
        srgb.set('val', '%02X%02X%02X' % (fill[0], fill[1], fill[2]))


# ════════════════════════════════════════════
# SLIDE 1: System Architecture Component Diagram
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg_white(slide)

# Title
txt(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
    "System Architecture — D-Bus IPC Component Diagram", 22, BLACK, True)

# Thin rule
box(slide, Inches(0.6), Inches(0.8), Inches(12.1), Pt(1), GRAY_30, GRAY_30, Pt(0.5))

# Subtitle
txt(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.3),
    "Manager(org.llm.Manager1)가 시스템 리소스를 모니터링/분석하여 LLM(Antigravity)에 D-Bus Signal을 단방향 전달한다.",
    10, GRAY_50)

# ── Component: Data Sources (leftmost) ──
DS_X = Inches(0.4)
DS_Y = Inches(1.7)
DS_W = Inches(2.2)
DS_H = Inches(5.0)

box(slide, DS_X, DS_Y, DS_W, DS_H, WHITE, BORDER)
txt(slide, DS_X, DS_Y + Pt(4), DS_W, Inches(0.3),
    "Data Sources", 11, GRAY_90, True, PP_ALIGN.CENTER)
box(slide, DS_X + Pt(2), DS_Y + Inches(0.35), DS_W - Pt(4), Pt(0.5), GRAY_30, GRAY_30, Pt(0.5))

sources = [
    ("/proc/meminfo", "Memory"),
    ("/proc/pressure/*", "PSI"),
    ("/proc/stat", "CPU"),
    ("/sys/class/devfreq/", "GPU"),
    ("thermal_zone*", "Thermal"),
    ("thermald (D-Bus)", "Thermal"),
    ("UPower (D-Bus)", "Battery"),
    ("/sys/class/power_supply/", "Power"),
]
for i, (path, _label) in enumerate(sources):
    y = DS_Y + Inches(0.55) + Inches(i * 0.52)
    rbox(slide, DS_X + Pt(6), y, DS_W - Pt(12), Inches(0.4), GRAY_08, GRAY_30, Pt(0.5))
    txt(slide, DS_X + Pt(10), y + Pt(1), DS_W - Pt(20), Inches(0.35),
        path, 8, GRAY_70, False, PP_ALIGN.LEFT, FONT_MONO)

# ── Arrow: Data Sources → Manager ──
A1_X1 = DS_X + DS_W
A1_X2 = Inches(3.1)
A1_Y = DS_Y + DS_H / 2
arrow_right(slide, A1_X1 + Pt(4), A1_Y, A1_X2 - Pt(4))

# ── Component: Manager (center-left) ──
MG_X = Inches(3.1)
MG_Y = Inches(1.7)
MG_W = Inches(3.3)
MG_H = Inches(5.0)

box(slide, MG_X, MG_Y, MG_W, MG_H, GRAY_05, GRAY_90, Pt(1.5))
txt(slide, MG_X, MG_Y + Pt(4), MG_W, Inches(0.35),
    "Manager", 14, BLACK, True, PP_ALIGN.CENTER)
txt(slide, MG_X, MG_Y + Inches(0.35), MG_W, Inches(0.25),
    "org.llm.Manager1", 9, GRAY_50, False, PP_ALIGN.CENTER, FONT_MONO)
box(slide, MG_X + Pt(4), MG_Y + Inches(0.6), MG_W - Pt(8), Pt(0.5), GRAY_30, GRAY_30, Pt(0.5))

steps = [
    ("1. Collect", "시스템 리소스 폴링 / D-Bus 수신"),
    ("2. Analyze", "추세 분석, 이동 평균, 변화율"),
    ("3. Predict", "임계 도달 예측, 쓰로틀링 예상"),
    ("4. Decide", "Level 결정 (히스테리시스 적용)"),
    ("5. Signal", "D-Bus Signal 발행"),
]
for i, (step, desc) in enumerate(steps):
    y = MG_Y + Inches(0.8) + Inches(i * 0.78)
    rbox(slide, MG_X + Pt(8), y, MG_W - Pt(16), Inches(0.62), WHITE, GRAY_30, Pt(0.5))
    txt(slide, MG_X + Pt(14), y + Pt(3), MG_W - Pt(28), Inches(0.22),
        step, 10, BLACK, True, PP_ALIGN.LEFT, FONT_BODY)
    txt(slide, MG_X + Pt(14), y + Pt(17), MG_W - Pt(28), Inches(0.22),
        desc, 8, GRAY_50, False, PP_ALIGN.LEFT, FONT_BODY)

# ── Arrow: Manager → D-Bus ──
A2_X1 = MG_X + MG_W
A2_X2 = Inches(6.9)
A2_Y = MG_Y + MG_H / 2
arrow_right(slide, A2_X1 + Pt(4), A2_Y, A2_X2 - Pt(4), GRAY_90, Pt(2))

# Arrow label
txt(slide, A2_X1 + Pt(2), A2_Y - Inches(0.3), Inches(0.6), Inches(0.25),
    "D-Bus", 8, GRAY_50, True, PP_ALIGN.CENTER)
txt(slide, A2_X1 + Pt(2), A2_Y - Inches(0.12), Inches(0.6), Inches(0.2),
    "Signal", 7, GRAY_50, False, PP_ALIGN.CENTER)

# ── Component: D-Bus Bus (center pipe) ──
BUS_X = Inches(6.9)
BUS_Y = Inches(1.7)
BUS_W = Inches(1.5)
BUS_H = Inches(5.0)

box(slide, BUS_X, BUS_Y, BUS_W, BUS_H, GRAY_08, GRAY_90, Pt(1))
txt(slide, BUS_X, BUS_Y + Pt(4), BUS_W, Inches(0.3),
    "System Bus", 10, GRAY_90, True, PP_ALIGN.CENTER)
box(slide, BUS_X + Pt(4), BUS_Y + Inches(0.35), BUS_W - Pt(8), Pt(0.5), GRAY_30, GRAY_30, Pt(0.5))

signals_list = [
    "MemoryPressure",
    "ComputeGuidance",
    "ThermalAlert",
    "EnergyConstraint",
]
for i, sig in enumerate(signals_list):
    y = BUS_Y + Inches(0.6) + Inches(i * 1.0)
    rbox(slide, BUS_X + Pt(6), y, BUS_W - Pt(12), Inches(0.7), WHITE, GRAY_50, Pt(0.5))
    txt(slide, BUS_X + Pt(8), y + Pt(4), BUS_W - Pt(16), Inches(0.55),
        sig, 8, BLACK, True, PP_ALIGN.CENTER, FONT_MONO)

# ── Arrow: D-Bus → LLM ──
A3_X1 = BUS_X + BUS_W
A3_X2 = Inches(8.9)
A3_Y = BUS_Y + BUS_H / 2
arrow_right(slide, A3_X1 + Pt(4), A3_Y, A3_X2 - Pt(4), GRAY_90, Pt(2))

# ── Component: LLM (right) ──
LLM_X = Inches(8.9)
LLM_Y = Inches(1.7)
LLM_W = Inches(4.0)
LLM_H = Inches(5.0)

box(slide, LLM_X, LLM_Y, LLM_W, LLM_H, GRAY_05, GRAY_90, Pt(1.5))
txt(slide, LLM_X, LLM_Y + Pt(4), LLM_W, Inches(0.35),
    "LLM Inference (Antigravity)", 14, BLACK, True, PP_ALIGN.CENTER)
txt(slide, LLM_X, LLM_Y + Inches(0.35), LLM_W, Inches(0.25),
    "Signal 수신 → 자율적 최적 동작", 9, GRAY_50, False, PP_ALIGN.CENTER)
box(slide, LLM_X + Pt(4), LLM_Y + Inches(0.6), LLM_W - Pt(8), Pt(0.5), GRAY_30, GRAY_30, Pt(0.5))

actions = [
    ("KV Cache Eviction", "메모리 확보를 위한 캐시 축소", "← MemoryPressure"),
    ("Backend Switch", "CPU ↔ GPU 동적 전환", "← ComputeGuidance"),
    ("Quality Adjustment", "Sampling, 토큰 수 제한", "← ThermalAlert"),
    ("Inference Control", "속도 제한, 일시 중단, 거부", "← EnergyConstraint"),
]
for i, (title, desc, trigger) in enumerate(actions):
    y = LLM_Y + Inches(0.8) + Inches(i * 0.97)
    rbox(slide, LLM_X + Pt(8), y, LLM_W - Pt(16), Inches(0.78), WHITE, GRAY_30, Pt(0.5))
    txt(slide, LLM_X + Pt(14), y + Pt(3), LLM_W - Pt(28), Inches(0.22),
        title, 10, BLACK, True, PP_ALIGN.LEFT, FONT_BODY)
    txt(slide, LLM_X + Pt(14), y + Pt(17), LLM_W - Pt(28), Inches(0.2),
        desc, 8, GRAY_50, False, PP_ALIGN.LEFT, FONT_BODY)
    txt(slide, LLM_X + Pt(14), y + Pt(30), LLM_W - Pt(28), Inches(0.2),
        trigger, 7, GRAY_30, False, PP_ALIGN.LEFT, FONT_MONO)

# Footer
txt(slide, Inches(0.6), Inches(7.0), Inches(4), Inches(0.3),
    "Bus: System Bus  |  Interface: org.llm.Manager1  |  Direction: Manager → LLM (단방향)",
    8, GRAY_50, False, PP_ALIGN.LEFT, FONT_MONO)


# ════════════════════════════════════════════
# SLIDE 2: D-Bus Signal Specification Table
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg_white(slide)

txt(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
    "D-Bus Signal Specification — org.llm.Manager1", 22, BLACK, True)
box(slide, Inches(0.6), Inches(0.8), Inches(12.1), Pt(1), GRAY_30, GRAY_30, Pt(0.5))
txt(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.3),
    "4개 시그널 상세 정의. 모든 시그널은 공통 level(normal / warning / critical / emergency) 인자를 포함한다.",
    10, GRAY_50)

# ── Main Signal Table ──
# Columns: Signal | Arguments | Type | Description | LLM Response
ROWS = 20  # header(1) + mem(4) + compute(6) + thermal(5) + energy(4)
COLS = 5
TBL_LEFT = Inches(0.5)
TBL_TOP = Inches(1.4)
TBL_W = Inches(12.3)
TBL_H = Inches(5.7)

tbl = add_table(slide, TBL_LEFT, TBL_TOP, TBL_W, TBL_H, ROWS, COLS)

# Column widths
col_widths = [Inches(2.0), Inches(2.3), Inches(0.6), Inches(4.0), Inches(3.4)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

# Header fill
HDR_FILL = (0x2A, 0x2A, 0x2A)  # dark gray
HDR_TEXT = WHITE
SECTION_FILL = (0xED, 0xED, 0xED)  # light gray section header
ROW_ALT = (0xF9, 0xF9, 0xF9)  # alternating row

headers = ["Signal", "Argument", "Type", "Description", "LLM Response"]
for c, h in enumerate(headers):
    style_cell(tbl.cell(0, c), h, 9, True, HDR_TEXT, HDR_FILL, PP_ALIGN.CENTER)

# ── Data rows ──
row = 1

# Helper to fill a row
def fill_row(r, sig, arg, typ, desc, response, section=False, alt=False):
    fill = SECTION_FILL if section else (ROW_ALT if alt else None)
    style_cell(tbl.cell(r, 0), sig, 9, section, BLACK if section else GRAY_90, fill, PP_ALIGN.LEFT, FONT_MONO if not section else FONT_BODY)
    style_cell(tbl.cell(r, 1), arg, 8, False, BLACK, fill, PP_ALIGN.LEFT, FONT_MONO)
    style_cell(tbl.cell(r, 2), typ, 8, False, GRAY_50, fill, PP_ALIGN.CENTER, FONT_MONO)
    style_cell(tbl.cell(r, 3), desc, 8, False, GRAY_70, fill, PP_ALIGN.LEFT)
    style_cell(tbl.cell(r, 4), response, 8, False, GRAY_70, fill, PP_ALIGN.LEFT)


# MemoryPressure section
fill_row(row, "MemoryPressure", "", "", "OOM 직전 경고. KV 캐시를 줄여야 하는 상황을 안내", "", section=True)
row += 1
fill_row(row, "", "level", "s", "normal | warning | critical | emergency", "warning: 보수적 eviction 준비")
row += 1
fill_row(row, "", "available_bytes", "t", "현재 사용 가능한 메모리 (bytes)", "critical: 적극적 KV 캐시 eviction", alt=True)
row += 1
fill_row(row, "", "reclaim_target_bytes", "t", "LLM이 확보해야 하는 메모리 양 (Manager 산출)", "emergency: 긴급 eviction + 신규 추론 거부")
row += 1

# ComputeGuidance section
fill_row(row, "ComputeGuidance", "", "", "CPU/GPU 병목 감지. 어떤 연산 HW를 사용할지 가이드", "", section=True)
row += 1
fill_row(row, "", "level", "s", "normal | warning | critical", "warning: 권장 백엔드로 전환 준비")
row += 1
fill_row(row, "", "recommended_backend", "s", "cpu | gpu | any (Manager 권장 백엔드)", "critical: 즉시 백엔드 전환 + KV 캐시 마이그레이션", alt=True)
row += 1
fill_row(row, "", "reason", "s", "cpu_bottleneck | gpu_bottleneck | cpu_available | gpu_available | both_loaded | balanced", "both_loaded: 추론 속도 제한 (throttle)")
row += 1
fill_row(row, "", "cpu_usage_pct", "d", "시스템 CPU 사용률 (0.0 ~ 100.0)", "", alt=True)
row += 1
fill_row(row, "", "gpu_usage_pct", "d", "GPU 사용률 (0.0 ~ 100.0)", "")
row += 1

# ThermalAlert section
fill_row(row, "ThermalAlert", "", "", "발열 쓰로틀링 발생 예측 및 실제 발생 상황을 안내", "", section=True)
row += 1
fill_row(row, "", "level", "s", "normal | warning | critical | emergency", "warning: GPU→CPU 전환 검토, 부하 경감 준비")
row += 1
fill_row(row, "", "temperature_mc", "i", "현재 온도 (millidegrees Celsius, 75000 = 75.0℃)", "critical: GPU 중지, 품질 하향", alt=True)
row += 1
fill_row(row, "", "throttling_active", "b", "HW 쓰로틀링이 현재 발생 중인지 여부", "emergency: 추론 일시 중단 (pause)")
row += 1
fill_row(row, "", "throttle_ratio", "d", "성능 비율 (1.0 = 정상, 0.5 = 절반 성능)", "", alt=True)
row += 1

# EnergyConstraint section
fill_row(row, "EnergyConstraint", "", "", "에너지 소모 경고 및 전력 예산 안내", "", section=True)
row += 1
fill_row(row, "", "level", "s", "normal | warning | critical | emergency", "warning: GPU 사용 자제, budget 내 동작")
row += 1
fill_row(row, "", "reason", "s", "battery_low | battery_critical | power_limit | thermal_power | charging | none", "critical: CPU only, 최소 품질, 토큰 제한", alt=True)
row += 1
fill_row(row, "", "power_budget_mw", "u", "허용 전력 예산 (milliwatts, 0 = 무제한)", "emergency: 신규 추론 거부, 진행 중 추론 종료")

# Footer note
txt(slide, Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.25),
    "Level 공통 규약: 상향 전이(악화) 즉시 반응 | 하향 전이(회복) 지연 허용 | normal 수신 시 모든 제약 해제 | Manager 미연결 시 normal 간주",
    8, GRAY_50, False, PP_ALIGN.LEFT)

# ════════════════════════════════════════════
# Save
# ════════════════════════════════════════════
output_path = "/home/go/Workspace/llm_rs2/docs/dbus_ipc_architecture.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
